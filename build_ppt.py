import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_overall_pipeline_pptx(output_path="figures_src/overall_pipeline_ch01.pptx"):
    prs = Presentation()
    # 16:9 Slide Dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6] # Blank

    # Color Palette: Academic Monochrome & Slate
    BG_COLOR = RGBColor(255, 255, 255)         # Pure White
    DARK_TEXT = RGBColor(26, 32, 44)           # Charcoal Black
    MUTED_TEXT = RGBColor(74, 85, 104)        # Slate Gray
    HEADER_BG = RGBColor(241, 245, 249)        # Very Light Blue-Gray
    CONTAINER_BG = RGBColor(248, 250, 252)     # Off-White / Pale Gray
    CONTAINER_BORDER = RGBColor(148, 163, 184) # Mid Slate Gray
    CARD_BG = RGBColor(255, 255, 255)          # White Card
    CARD_BORDER = RGBColor(203, 213, 225)      # Light Slate Gray
    ACCENT_BORDER = RGBColor(71, 85, 105)      # Dark Slate
    ACCENT_BG = RGBColor(226, 232, 240)        # Accent Gray

    FONT_HEADING = "Arial"
    FONT_BODY = "Microsoft YaHei" # Chinese compatibility

    # =========================================================================
    # SLIDE 1: Master Overall Pipeline Diagram (主框架整体流程图)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_slide_layout)

    # 1. Slide Title Header
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0
    p1 = tf1.paragraphs[0]
    p1.text = "专著整体框架与章节映射流程图 (Overall Framework & Chapter Structure)"
    p1.font.name = FONT_HEADING
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = DARK_TEXT

    p1_sub = tf1.add_paragraph()
    p1_sub.text = "书名: AI-Driven UAV Building Inspection: An End-to-End Framework from Autonomous Flight to Digital Twins (第1章导论映射)"
    p1_sub.font.name = FONT_BODY
    p1_sub.font.size = Pt(12)
    p1_sub.font.color.rgb = MUTED_TEXT

    # 2. Main Container Box (The 4 Horizontal Phases)
    phases = [
        {
            "phase": "阶段一: 自主飞行与数据采集",
            "en": "Phase I: Autonomous Flight & Acquisition",
            "bg": RGBColor(248, 250, 252),
            "border": RGBColor(148, 163, 184),
            "chapters": [
                {
                    "num": "第 2 章",
                    "title": "任务规划 (Task Planning)",
                    "desc": "立面巡检覆盖规划 • ROI区域定义 • 多无人机路径生成"
                },
                {
                    "num": "第 3 章",
                    "title": "运动规划与执行 (Motion Planning)",
                    "desc": "实时避障 • 轨迹优化与控制 • 风场适应与自主飞行"
                }
            ]
        },
        {
            "phase": "阶段二: 空间重建与数据表达",
            "en": "Phase II: Spatial Reconstruction & Representation",
            "bg": RGBColor(248, 250, 252),
            "border": RGBColor(148, 163, 184),
            "chapters": [
                {
                    "num": "第 4 章",
                    "title": "3D重建与资产映射 (3D Reconstruction)",
                    "desc": "摄影测量 • 点云与网格建模 • BIM/坐标系精准对齐"
                },
                {
                    "num": "第 5 章",
                    "title": "巡检数据集与多模态数据 (Datasets)",
                    "desc": "RGB/热成像数据采集 • 缺陷标注规范 • 基准集构建"
                }
            ]
        },
        {
            "phase": "阶段三: AI 缺陷感知与分析",
            "en": "Phase III: AI Defect Perception & Analysis",
            "bg": RGBColor(248, 250, 252),
            "border": RGBColor(148, 163, 184),
            "chapters": [
                {
                    "num": "第 6 章",
                    "title": "缺陷检测 AI 模型 (AI Defect Models)",
                    "desc": "深度学习检测器 • 实例分割 • 缺陷面积与程度量化"
                }
            ]
        },
        {
            "phase": "阶段四: 数字孪生与 AI 决策支持",
            "en": "Phase IV: Digital Twin & AI Decision Support",
            "bg": RGBColor(241, 245, 249),
            "border": RGBColor(71, 85, 105),
            "chapters": [
                {
                    "num": "第 7 章",
                    "title": "GIS 与 GeoBIM 集成 (GeoBIM Integration)",
                    "desc": "空间数据库映射 • GeoBIM Schema • 资产护照挂载"
                },
                {
                    "num": "第 8 章",
                    "title": "数字孪生与大模型 (DT & LLM Reasoning)",
                    "desc": "LLM/VLM推理 • 检索增强生成(RAG) • 预测性维护决策"
                },
                {
                    "num": "第 9 章",
                    "title": "总结与展望 (Conclusion & Future)",
                    "desc": "现场部署经验审计 • 开放挑战与下一代智能运维展望"
                }
            ]
        }
    ]

    # Layout dimensions for 4 columns
    start_x = Inches(0.8)
    col_width = Inches(2.7)
    gap = Inches(0.3)
    start_y = Inches(1.5)
    container_height = Inches(5.1)

    for i, phase_data in enumerate(phases):
        col_x = start_x + i * (col_width + gap)

        # Draw Phase Container Background Shape
        rect = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x, start_y, col_width, container_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = phase_data["bg"]
        rect.line.color.rgb = phase_data["border"]
        rect.line.width = Pt(1.5)

        # Header Box for Phase
        header_shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x, start_y, col_width, Inches(0.65))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = phase_data["border"]
        header_shape.line.fill.background()

        tf_h = header_shape.text_frame
        tf_h.word_wrap = True
        tf_h.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_h1 = tf_h.paragraphs[0]
        p_h1.text = phase_data["phase"]
        p_h1.font.name = FONT_BODY
        p_h1.font.size = Pt(11)
        p_h1.font.bold = True
        p_h1.font.color.rgb = RGBColor(255, 255, 255)
        p_h1.alignment = PP_ALIGN.CENTER

        # Draw Cards inside Phase
        card_y = start_y + Inches(0.75)
        num_chaps = len(phase_data["chapters"])
        available_h = container_height - Inches(0.85)
        card_h = (available_h - gap * (num_chaps - 1)) / num_chaps if num_chaps > 0 else Inches(1.2)

        for ch in phase_data["chapters"]:
            card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col_x + Inches(0.12), card_y, col_width - Inches(0.24), card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = CARD_BG
            card.line.color.rgb = CARD_BORDER
            card.line.width = Pt(1.0)

            tf_c = card.text_frame
            tf_c.word_wrap = True
            tf_c.margin_left = Inches(0.1)
            tf_c.margin_right = Inches(0.1)
            tf_c.margin_top = Inches(0.08)
            tf_c.margin_bottom = Inches(0.08)

            p_c1 = tf_c.paragraphs[0]
            p_c1.text = f"【{ch['num']}】{ch['title']}"
            p_c1.font.name = FONT_BODY
            p_c1.font.size = Pt(11)
            p_c1.font.bold = True
            p_c1.font.color.rgb = DARK_TEXT

            p_c2 = tf_c.add_paragraph()
            p_c2.text = ch["desc"]
            p_c2.font.name = FONT_BODY
            p_c2.font.size = Pt(9.5)
            p_c2.font.color.rgb = MUTED_TEXT

            card_y += card_h + Inches(0.1)

        # Draw Inter-Phase Connecting Arrows (between column i and i+1)
        if i < len(phases) - 1:
            arrow_x = col_x + col_width + Inches(0.05)
            arrow_y = start_y + container_height / 2 - Inches(0.15)
            arrow = slide1.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, Inches(0.2), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED_TEXT
            arrow.line.fill.background()

    # Footer Note
    footer_box = slide1.shapes.add_textbox(Inches(0.8), Inches(6.75), Inches(11.733), Inches(0.4))
    tf_f = footer_box.text_frame
    p_f = tf_f.paragraphs[0]
    p_f.text = "注：从第2/3章自主飞行与采集，到第4/5章空间重建，第6章缺陷识别，最终至第7/8/9章数字孪生与大模型决策，形成完整的端到端智能运维闭环。"
    p_f.font.name = FONT_BODY
    p_f.font.size = Pt(10)
    p_f.font.italic = True
    p_f.font.color.rgb = MUTED_TEXT

    # =========================================================================
    # SLIDE 2: Information Flow & Inter-Chapter Interfaces (章节信息流转与接口)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_slide_layout)

    # Title
    title_box2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
    tf2 = title_box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "全书章节数据流转与功能接口 (Data Flow & Inter-Chapter Interfaces)"
    p2.font.name = FONT_HEADING
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = DARK_TEXT

    p2_sub = tf2.add_paragraph()
    p2_sub.text = "构建从底层物理采样到高层大模型推理的完整数据链条 (Data Chain from Physical Sensing to Cognitive LLM)"
    p2_sub.font.name = FONT_BODY
    p2_sub.font.size = Pt(12)
    p2_sub.font.color.rgb = MUTED_TEXT

    # Matrix Table for Data Flow
    rows = 6
    cols = 4
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(11.733)
    height = Inches(5.2)

    table_shape = slide2.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Column Widths
    table.columns[0].width = Inches(2.2) # Source Chapter
    table.columns[1].width = Inches(2.8) # Output Data / Asset
    table.columns[2].width = Inches(3.9) # Functional Operations
    table.columns[3].width = Inches(2.833) # Downstream Target Chapter

    headers = ["阶段 / 章节 (Phase & Chapter)", "输出数据/资产 (Data & Asset)", "核心处理/功能 (Key Operation)", "下游对接章节 (Target Chapter)"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(226, 232, 240)
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER

    matrix_data = [
        ("第 2 & 3 章: 任务与运动规划\n(Flight & Motion Planning)", "三维无碰撞巡检轨迹\n多模态影像数据集 (RGB/Thermal)", "覆盖率优化路径生成、实时避障与风场扰动控制", "第 4 & 5 章: 重建与数据集"),
        ("第 4 & 5 章: 3D重建与数据集\n(Reconstruction & Datasets)", "高精度点云、网格模型(Mesh)\n标注对齐的缺陷图像集", "倾斜摄影测量、BIM坐标系精准对齐、标定与分割", "第 6 章: AI缺陷检测\n第 7 章: GIS/GeoBIM"),
        ("第 6 章: 缺陷检测 AI 模型\n(AI Defect Perception)", "缺陷 bounding boxes/掩码\n构件级缺陷类型、面积与定量指标", "深度学习目标检测、实例分割、缺陷严重程度量化", "第 7 章: GeoBIM 挂载\n第 8 章: 数字孪生推理"),
        ("第 7 章: GIS 与 GeoBIM 集成\n(GeoBIM & Asset Passports)", "挂载缺陷的 GeoBIM 资产护照\n空间拓扑关系数据库 (Spatial Graph)", "空间关联分析、Asset-ID 语义绑定、历史维保数据库归档", "第 8 章: 大模型决策支持"),
        ("第 8 & 9 章: 数字孪生与大模型\n(DT & LLM Decision Support)", "可审计巡检报告、维保建议工单\n闭环复巡指令 (Re-inspection Task)", "检索增强生成 (RAG)、历史记录查询、预测性维护决策", "反馈至 第 2 章 闭环复巡\n及现场运维工程应用")
    ]

    for row_idx, row_content in enumerate(matrix_data, start=1):
        for col_idx, cell_text in enumerate(row_content):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_idx % 2 == 1 else RGBColor(248, 250, 252)
            p = cell.text_frame.paragraphs[0]
            p.text = cell_text
            p.font.name = FONT_BODY
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_TEXT

    # =========================================================================
    # SLIDE 3: System Architecture Layer View (物理-空间-感知-认知 四层架构视图)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_slide_layout)

    # Title
    title_box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.9))
    tf3 = title_box3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "建筑巡检端到端四层系统架构视图 (Four-Layer System Architecture)"
    p3.font.name = FONT_HEADING
    p3.font.size = Pt(22)
    p3.font.bold = True
    p3.font.color.rgb = DARK_TEXT

    p3_sub = tf3.add_paragraph()
    p3_sub.text = "物理域、空间域、感知域与认知域的解耦与协同 (Decoupled Co-Design from Flight to Cognitive Reasoning)"
    p3_sub.font.name = FONT_BODY
    p3_sub.font.size = Pt(12)
    p3_sub.font.color.rgb = MUTED_TEXT

    # 4 Horizontal Stack Layers
    layers_data = [
        {
            "layer_name": "L4: 认知与决策层 (Cognitive & Decision Domain)",
            "chaps": "对应章节: 第 8 章 (DT + LLM) & 第 9 章 (总结与展望)",
            "content": "即用型大语言模型 (LLM/VLM) • 检索增强生成 (RAG) • 历史维修记录查询 • 预测性维护决策 • 报告自动生成与工单下发",
            "bg": RGBColor(241, 245, 249),
            "border": RGBColor(71, 85, 105)
        },
        {
            "layer_name": "L3: 感知与资产层 (Perception & Asset Management Domain)",
            "chaps": "对应章节: 第 6 章 (AI检测模型) & 第 7 章 (GIS & GeoBIM集成)",
            "content": "深度学习缺陷检测/实例分割 • 图像语义分析 • GeoBIM 空间数据库 • 资产护照 (Asset Passport) • 拓扑关系链锚定",
            "bg": RGBColor(248, 250, 252),
            "border": RGBColor(148, 163, 184)
        },
        {
            "layer_name": "L2: 空间与表达层 (Spatial & Data Representation Domain)",
            "chaps": "对应章节: 第 4 章 (3D重建) & 第 5 章 (巡检数据集)",
            "content": "无人机倾斜摄影测量 • 点云与三维网格建模 (Mesh) • BIM/地理坐标对齐 • 多模态 (RGB/热成像) 数据集构建与标定",
            "bg": RGBColor(248, 250, 252),
            "border": RGBColor(148, 163, 184)
        },
        {
            "layer_name": "L1: 物理与采集层 (Physical & Flight Domain)",
            "chaps": "对应章节: 第 2 章 (任务规划) & 第 3 章 (运动规划与执行)",
            "content": "立面巡检覆盖路径规划 • 实时避障与安全边界 • 轨迹优化控制 • 风场扰动控制 • 无人机平台与多传感器硬件执行",
            "bg": RGBColor(248, 250, 252),
            "border": RGBColor(148, 163, 184)
        }
    ]

    stack_y = Inches(1.5)
    stack_h = Inches(1.15)
    stack_gap = Inches(0.15)

    for idx, ldata in enumerate(layers_data):
        curr_y = stack_y + idx * (stack_h + stack_gap)

        # Background Container
        s_shape = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), curr_y, Inches(11.733), stack_h)
        s_shape.fill.solid()
        s_shape.fill.fore_color.rgb = ldata["bg"]
        s_shape.line.color.rgb = ldata["border"]
        s_shape.line.width = Pt(1.5)

        # Text Frame
        tf_s = s_shape.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = Inches(0.2)
        tf_s.margin_right = Inches(0.2)
        tf_s.margin_top = Inches(0.1)
        tf_s.margin_bottom = Inches(0.1)

        p_l1 = tf_s.paragraphs[0]
        p_l1.text = f"{ldata['layer_name']}   │   {ldata['chaps']}"
        p_l1.font.name = FONT_BODY
        p_l1.font.size = Pt(12)
        p_l1.font.bold = True
        p_l1.font.color.rgb = DARK_TEXT

        p_l2 = tf_s.add_paragraph()
        p_l2.text = f"核心功能组件: {ldata['content']}"
        p_l2.font.name = FONT_BODY
        p_l2.font.size = Pt(10.5)
        p_l2.font.color.rgb = MUTED_TEXT

        # Upward/Downward Inter-Layer Arrows
        if idx < len(layers_data) - 1:
            arrow_y = curr_y + stack_h + Inches(0.02)
            arrow_x = Inches(6.4)
            arrow = slide3.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, arrow_x, arrow_y, Inches(0.25), Inches(0.11))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED_TEXT
            arrow.line.fill.background()

    # Ensure output folder exists
    dirname = os.path.dirname(output_path)
    if dirname:
        os.makedirs(dirname, exist_ok=True)
    prs.save(output_path)
    print(f"Successfully generated presentation at: {output_path}")

if __name__ == "__main__":
    create_overall_pipeline_pptx("figures_src/overall_pipeline_ch01.pptx")
    # Also save to root workspace for easy access
    create_overall_pipeline_pptx("Chapter01_Overall_Pipeline.pptx")
